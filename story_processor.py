import os
import json
import argparse
import re
from pathlib import Path
from tqdm import tqdm
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import spacy
import tiktoken

# Load the spaCy NLP model for sentence and paragraph detection
nlp = spacy.load("en_core_web_sm")

# Initialize tokenizer for counting tokens (using cl100k_base tokenizer used by many models)
tokenizer = tiktoken.get_encoding("cl100k_base")

def extract_text_from_docx(file_path):
    """Extract text from a .docx file, preserving paragraph breaks."""
    try:
        doc = Document(file_path)
        paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        return paragraphs
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return []

def extract_text_from_pptx(file_path):
    """Extract text from a .pptx file."""
    try:
        prs = Presentation(file_path)
        paragraphs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    paragraphs.append(shape.text.strip())
        return paragraphs
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return []

def extract_text_from_html(file_path):
    """Extract text from an .html file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            # Try to extract paragraphs with proper tag separation
            paragraphs = []
            for p in soup.find_all(['p', 'div', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                if p.text.strip():
                    paragraphs.append(p.text.strip())
            # If no paragraphs found, fall back to simple text extraction
            if not paragraphs:
                text = soup.get_text(separator="\n").strip()
                paragraphs = [p for p in text.split("\n") if p.strip()]
            return paragraphs
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return []

def extract_text_from_txt(file_path):
    """Extract text from a .txt file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            # Split by double newlines to preserve paragraph structure
            content = f.read().strip()
            paragraphs = [p.strip() for p in re.split(r'\n\s*\n', content) if p.strip()]
            return paragraphs
    except Exception as e:
        print(f"Error processing {file_path}: {e}")
        return []

def split_into_sentences(text):
    """Split text into sentences using spaCy."""
    doc = nlp(text)
    return [sent.text.strip() for sent in doc.sents if sent.text.strip()]

def count_tokens(text):
    """Count the number of tokens in a text using tiktoken."""
    tokens = tokenizer.encode(text)
    return len(tokens)

def create_chunks_by_tokens(paragraphs, max_tokens=2048, overlap_paragraphs=1):
    """Create chunks of the story based on token count with paragraph overlap."""
    chunks = []
    current_chunk = []
    current_token_count = 0
    
    for i, para in enumerate(paragraphs):
        para_tokens = count_tokens(para)
        
        # If this paragraph alone exceeds max tokens, split it into sentences
        if para_tokens > max_tokens:
            sentences = split_into_sentences(para)
            sentence_chunk = []
            sentence_token_count = 0
            
            for sentence in sentences:
                sentence_tokens = count_tokens(sentence)
                if sentence_token_count + sentence_tokens <= max_tokens:
                    sentence_chunk.append(sentence)
                    sentence_token_count += sentence_tokens
                else:
                    if sentence_chunk:  # Save current sentence chunk
                        chunks.append(" ".join(sentence_chunk))
                    sentence_chunk = [sentence]  # Start new chunk with this sentence
                    sentence_token_count = sentence_tokens
            
            if sentence_chunk:  # Save the last sentence chunk if it exists
                chunks.append(" ".join(sentence_chunk))
            continue
        
        # If adding this paragraph exceeds the limit, save current chunk and start new one
        if current_token_count + para_tokens > max_tokens and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            # Keep overlap paragraphs for context continuity
            if overlap_paragraphs > 0 and len(current_chunk) >= overlap_paragraphs:
                current_chunk = current_chunk[-overlap_paragraphs:]
                current_token_count = count_tokens("\n\n".join(current_chunk))
            else:
                current_chunk = []
                current_token_count = 0
        
        current_chunk.append(para)
        current_token_count += para_tokens
    
    # Add the last chunk if there's anything left
    if current_chunk:
        chunks.append("\n\n".join(current_chunk))
    
    return chunks

def create_training_examples(chunks, format_type="alpaca"):
    """Create training examples from chunks based on the specified format."""
    examples = []
    
    for i, chunk in enumerate(chunks):
        if format_type == "alpaca":
            # Alpaca instruction format
            instruction = "Continue this story in the same style and tone:"
            context = chunk if i == 0 else chunks[i-1]  # Use previous chunk as context when available
            continuation = chunk
            
            example = {
                "instruction": instruction,
                "input": context,
                "output": continuation
            }
        
        elif format_type == "chatml":
            # ChatML format (for chat models)
            messages = [
                {"role": "system", "content": "You are a creative storyteller. Continue the story when asked."},
                {"role": "user", "content": f"Continue this story: {chunk if i == 0 else chunks[i-1]}"},
                {"role": "assistant", "content": chunk}
            ]
            example = {"messages": messages}
            
        elif format_type == "completion":
            # Simple prompt-completion format (original style)
            sentences = split_into_sentences(chunk)
            prompt_end = min(3, len(sentences))
            
            example = {
                "prompt": " ".join(sentences[:prompt_end]),
                "completion": chunk
            }
        
        examples.append(example)
    
    return examples

def process_folder(folder_path, output_jsonl, max_tokens=2048, overlap_paragraphs=1, format_type="alpaca"):
    """Process all supported files in the folder and create a JSONL file with chunks."""
    if not os.path.isdir(folder_path):
        print(f"Error: {folder_path} is not a valid directory.")
        return
    
    file_paths = []
    for root, _, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            if file.lower().endswith(('.docx', '.pptx', '.html', '.txt')):
                file_paths.append(file_path)
    
    total_examples = 0
    with open(output_jsonl, "w", encoding="utf-8") as f:
        for file_path in tqdm(file_paths, desc="Processing files"):
            file_name = os.path.basename(file_path)
            
            if file_name.lower().endswith(".docx"):
                paragraphs = extract_text_from_docx(file_path)
            elif file_name.lower().endswith(".pptx"):
                paragraphs = extract_text_from_pptx(file_path)
            elif file_name.lower().endswith(".html"):
                paragraphs = extract_text_from_html(file_path)
            elif file_name.lower().endswith(".txt"):
                paragraphs = extract_text_from_txt(file_path)
            else:
                continue
            
            if not paragraphs:
                print(f"No text extracted from: {file_name}")
                continue
            
            # Create chunks based on token count
            chunks = create_chunks_by_tokens(paragraphs, max_tokens, overlap_paragraphs)
            
            # Skip if no chunks were created
            if not chunks:
                print(f"No valid chunks created for: {file_name}")
                continue
            
            # Create examples from chunks
            examples = create_training_examples(chunks, format_type)
            
            # Write examples to JSONL file
            for example in examples:
                f.write(json.dumps(example) + "\n")
                total_examples += 1
            
            print(f"Processed {file_name}: Created {len(examples)} examples")
    
    return total_examples

def main():
    parser = argparse.ArgumentParser(description="Process story files for LLM fine-tuning with Unsloth")
    parser.add_argument("--folder", type=str, required=True, help="Path to the folder containing story files")
    parser.add_argument("--output", type=str, default="training_data.jsonl", help="Output JSONL file path")
    parser.add_argument("--max-tokens", type=int, default=4096, help="Maximum tokens per chunk (default: 4096)")
    parser.add_argument("--overlap", type=int, default=1, help="Number of paragraphs to overlap between chunks")
    parser.add_argument("--format", type=str, choices=["alpaca", "chatml", "completion"], default="alpaca", 
                        help="Output format type (alpaca, chatml, or completion)")
    
    args = parser.parse_args()
    
    print(f"Starting processing of folder: {args.folder}")
    print(f"Using {args.format} format with max {args.max_tokens} tokens per chunk")
    
    total_examples = process_folder(
        args.folder, 
        args.output, 
        max_tokens=args.max_tokens,
        overlap_paragraphs=args.overlap,
        format_type=args.format
    )
    
    print(f"Training data saved to: {args.output}")
    print(f"Total training examples created: {total_examples}")

if __name__ == "__main__":
    main()